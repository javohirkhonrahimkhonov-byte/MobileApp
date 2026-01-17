
import os
import io
import logging
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

def extract_text_from_file(file_path: str, file_ext: str) -> str:
    """
    Turli fayl formatlaridan (PDF, DOCX, PPTX, TXT) matnni ajratib oladi.
    """
    full_text = ""
    file_ext = file_ext.lower().replace(".", "")

    try:
        if file_ext == "pdf":
            reader = PdfReader(file_path)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    full_text += text + "\n"
        
        elif file_ext == "docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                full_text += para.text + "\n"
        
        elif file_ext == "pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
        
        elif file_ext == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                full_text = f.read()

        else:
            return "Kechirasiz, bu fayl formati hozircha qo'llab-quvvatlanmaydi."

    except Exception as e:
        logger.error(f"Faylni o'qishda xatolik: {e}")
        return "Faylni o'qishda xatolik yuz berdi. Fayl shikastlangan bo'lishi mumkin."

    return full_text.strip()
